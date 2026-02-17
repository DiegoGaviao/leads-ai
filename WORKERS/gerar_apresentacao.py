
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    # 1. Configurar Slides e Estilo Clean
    prs = Presentation()
    
    # Função auxiliar para criar slide de Título
    def add_title_slide(title, subtitle):
        slide_layout = prs.slide_layouts[0] # 0 = Title Slide
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        subtitle_placeholder = slide.placeholders[1]
        
        title_placeholder.text = title
        subtitle_placeholder.text = subtitle

    # Função auxiliar para criar slide de Conteúdo
    def add_content_slide(title, content_list):
        slide_layout = prs.slide_layouts[1] # 1 = Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Título
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Conteúdo
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        
        for item in content_list:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(20)
            p.level = 0

    # --- SLIDE 1: CAPA ---
    add_title_slide("Mentoria Emocional para Gestantes", "Um enxoval invisível. Mas essencial.\nProposta por Karina Botti")

    # --- SLIDE 2: INTRODUÇÃO (A Dor) ---
    add_content_slide("O Enxoval Emocional", [
        "Quando uma mulher descobre a gravidez, ela monta o enxoval do bebê.",
        "Mas quase ninguém ensina a montar o Enxoval da Alma.",
        "Um espaço para acolher medo, culpa, expectativas e luto.",
        "Ordenar por dentro... antes de ter que segurar tudo por fora."
    ])

    # --- SLIDE 3: JUSTIFICATIVA ---
    add_content_slide("Por que essa Mentoria Precisa Existir?", [
        "O programa atual já cuida do físico (corpo, bebê).",
        "Mas e o coração dessa mulher?",
        "Gestação e puerpério são portais que abrem feridas e memórias.",
        "Sem um espaço seguro, o nascimento pode virar colapso.",
        "A mentoria torna o cuidado realmente INTEGRAL."
    ])

    # --- SLIDE 4: SUGESTÕES DE NOME ---
    add_content_slide("Identidade da Mentoria", [
        "Opção 1: Enxoval da Alma",
        "Opção 2: Antes do Colo, o Centro",
        "Opção 3: Raízes Maternas",
        "Opção 4: Ouvindo o Coração da Mãe",
        "Opção 5: Nascer Mãe"
    ])

    # --- SLIDE 5: OBJETIVOS ---
    add_content_slide("Objetivos do Programa", [
        "Acolher emoções (medo, culpa) da gestação/pós-parto.",
        "Prevenir crises (Burnout Materno, DPP, Ansiedade).",
        "Reforçar a identidade da mulher além da maternidade.",
        "Preparar o emocional para a mudança conjugal.",
        "Ensinar práticas de autoacolhimento e fé."
    ])

    # --- SLIDE 6: FORMATO ---
    add_content_slide("Estrutura da Mentoria", [
        "Formato: Grupo Online (com encontro presencial opcional).",
        "Duração: 6 Encontros (Semanais ou Quinzenais).",
        "Suporte: Grupo de WhatsApp para direção leve.",
        "Material: Exercícios práticos + PDFs de rituais.",
        "Encerramento: Roda presencial com simbologias (flores, oração)."
    ])

    # --- SLIDE 7: TEMAS DOS ENCONTROS ---
    add_content_slide("Jornada de 6 Encontros", [
        "1. O que está nascendo em mim? (Nova Identidade)",
        "2. Não sou só mãe (Valor da Mulher Inteira)",
        "3. O que eu esperava e o que é (Frustrações e Culpa)",
        "4. O amor que também assusta (Medos e Instabilidades)",
        "5. Cuidar de mim não é luxo (Rotina com Sentido)",
        "6. Eu não volto a ser a mesma. Mas posso voltar a mim."
    ])

    # --- SLIDE 8: PÚBLICO ALVO ---
    add_content_slide("Para Quem É?", [
        "Gestantes a partir da 16ª semana.",
        "Puérperas até 6 meses.",
        "Mulheres atendidas pelo programa multidisciplinar.",
        "Quem sente: medo, exaustão, insegurança, perda de si."
    ])

    # --- SLIDE 9: BENEFÍCIOS (MULHER & PROGRAMA) ---
    add_content_slide("Ganhos Reais", [
        "Para a Mulher: Clareza, Leveza, Conexão profunda com bebê/parceiro.",
        "Para o Programa do Obstetra: Cuidado Integral, Prevenção de crises.",
        "Diferencial: Aumento da fidelização e recomendação.",
        "Ética: Alinhamento com visão de saúde total."
    ])

    # --- SLIDE 10: FECHAMENTO ---
    add_title_slide("Um Colo para a Alma", 
    "\"Quando ela é cuidada, ela cuida com presença e verdade.\"\n\nEsse é o enxoval que ninguém vê. Mas que faz toda a diferença.\n🌹")

    # Salvar
    prs.save("Mentoria_Enxoval_Emocional.pptx")
    print("Apresentação criada com sucesso: Mentoria_Enxoval_Emocional.pptx")

if __name__ == "__main__":
    create_presentation()
